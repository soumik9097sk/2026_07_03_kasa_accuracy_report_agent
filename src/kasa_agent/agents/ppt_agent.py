from __future__ import annotations

import os
import re
from datetime import datetime
from typing import Any, Dict, Optional

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from kasa_agent.config import OUTPUT_DIR
from kasa_agent.services.chart_services import generate_actual_vs_forecast_chart

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
BODY_FONT_SIZE = Pt(12)


def _flatten_paragraph(paragraph) -> None:
    """
    Force zero left margin/hanging-indent and no auto-bullet on a paragraph.
    `level` alone doesn't guarantee this in a plain textbox -- the theme's
    default list style can still apply indent space -- so this is set
    directly on the paragraph XML instead.
    """
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        existing = pPr.find(qn(tag))
        if existing is not None:
            pPr.remove(existing)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _default_output_path(report: Dict[str, Any]) -> str:
    """One file per query, saved to output/ with a timestamp so repeat runs never collide."""
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    entity_slug = re.sub(r"[^A-Za-z0-9]+", "_", report["entity_value"]).strip("_")
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{entity_slug}_{report['year']}_{report['month']:02d}_{timestamp}.pptx"
    return os.path.join(OUTPUT_DIR, filename)


def _build_heading(report: Dict[str, Any]) -> str:
    entity = report["entity_value"].replace(" / ", " ").upper()
    return f"{entity} {report['year']} MONTH {report['month']} FINDINGS"


def _add_title(slide, heading: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.25), SLIDE_WIDTH - Inches(0.8), Inches(0.8)
    )
    tf = box.text_frame
    tf.text = heading
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True


def _add_findings_bullets(slide, report: Dict[str, Any]) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.3), Inches(6.1), SLIDE_HEIGHT - Inches(1.7)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for idx, kf in enumerate(report["key_findings"]):
        header_p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        header_p.text = f"{kf['key']} ({kf['bias'].replace('_', ' ')})"
        header_p.font.size = BODY_FONT_SIZE
        header_p.font.bold = True
        _flatten_paragraph(header_p)

        yoy = kf.get("yoy_sales") or {}
        if yoy.get("last_year_actual") is not None:
            change_pct = yoy.get("yoy_change_pct")
            change_text = f"{change_pct:+.1f}%" if change_pct is not None else "N/A"
            baseline_p = tf.add_paragraph()
            baseline_p.text = (
                f"- Baseline: {yoy['this_year_actual']:,.0f} units this year vs "
                f"{yoy['last_year_actual']:,.0f} units last year ({change_text})"
            )
            baseline_p.font.size = BODY_FONT_SIZE
            baseline_p.font.italic = True
            _flatten_paragraph(baseline_p)

        for reason in kf["reasons"]:
            reason_p = tf.add_paragraph()
            reason_p.text = f"- {reason}"
            reason_p.font.size = BODY_FONT_SIZE
            _flatten_paragraph(reason_p)


def _add_chart(slide, chart_path: str) -> None:
    slide.shapes.add_picture(
        chart_path, Inches(6.8), Inches(1.3), width=Inches(6.1)
    )


def generate_findings_slide(
    report: Dict[str, Any],
    chart_path: Optional[str] = None,
    output_path: Optional[str] = None,
) -> str:
    """
    Build a single findings slide from a generate_accuracy_report() result:
    heading (e.g. "AMAZON STAND MIXERS 2026 MONTH 5 FINDINGS"), left half
    bullet points of each key's bias + reasons, right half an actual vs
    forecast bar chart for those keys. Saves the .pptx and returns its path.

    Each call produces exactly one slide in a new file -- by default a
    timestamped path under output/, so repeat runs for the same query never
    overwrite or append to a previous one.
    """
    output_path = output_path or _default_output_path(report)
    chart_path = chart_path or os.path.splitext(output_path)[0] + "_chart.png"
    chart_path = generate_actual_vs_forecast_chart(
        report["key_findings"], title="Actual vs Forecast", save_path=chart_path
    )

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    _add_title(slide, _build_heading(report))
    _add_findings_bullets(slide, report)
    _add_chart(slide, chart_path)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)

    prs.save(output_path)
    return output_path
